
import pandas as pd 
from sklearn.tree import DecisionTreeClassifier, plot_tree
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
import os

data={"Product type":[250,100],
"Cost Price":[375.00,165.00],
"Selling Price":[480.00,200.00],
"Full Quantity":[60,150],
"Profit":[105.00,35.00],
"QuantitySold":[15,20]}

df=pd.DataFrame(data)


print(data)
x=df[["Cost Price", "Selling Price", "Full Quantity", "Profit"]]
y=df["QuantitySold"]

model=DecisionTreeClassifier(max_depth=3, min_samples_leaf=1, random_state=42)
model.fit(x,y)

#visulaize the tree
plt.figure(figsize=(12,8))
plot_tree(
model,
feature_names=["Cost Price", "Selling Price", "Full Quantity", "Profit"],
class_names=[str(c) for c in model.classes_],
filled=True,
rounded=True,
fontsize=10)

plt.show()

tree_path="trr.png"
plt.savefig(tree_path)

plt.close()
#Feature Importance Plot
feature_importances_path="feature.png"
plt.barh(model.feature_names_in_,model.feature_importances_,color='blue')
plt.xlabel("feature importance")
plt.title("feature importance from decision tree")
plt.tight_layout()
plt.savefig(feature_importances_path)
plt.show()
plt.close()


#Product Type Verses QuantitySold
path="feature.png"
labels=[f"{p}g" for p in df["Product type"]]
labels2=[int(a) for a in df["QuantitySold"]]
plt.bar(labels,labels2,color='blue')

plt.xlabel("Product type")
plt.ylabel("Quantity Sold")
plt.ylim(0,20)

plt.title("Product Type Verses Quantity Sold")
plt.tight_layout()
plt.savefig(path)
plt.show()
plt.close()


new_data=pd.DataFrame({
    "Cost Price": [200, 300, 400],
    "Selling Price": [250, 450, 500],
    "Full Quantity": [100, 80, 50],
    "Profit": [50, 150, 100]
})

#Make Predctions
predictions=model.predict(new_data)
print(new_data)
print("Predictions:", predictions)

# Ai recommendation function
def teaAdvisor(cost,selling,profit,predict_sale):
    if profit<50:
        return "Low profit margin"
    elif profit>150:
        return "Promote this product to improve sales"
    elif predict_sale<10:
        return "low sales.consider discounts"
    elif predict_sale>50:
        return "High sales.Consider increasing price"
    else:
        return "Balanced Sales and profits"
    
    #Generate advice
advice_results=[]

for i in range(len(new_data)):
    cost=new_data.iloc[i]["Cost Price"]
    selling=new_data.iloc[i]["Selling Price"]
    profit=new_data.iloc[i]["Profit"]
    predict_sale=predictions[i]

    advice=teaAdvisor(cost,selling,profit,predict_sale)
    advice_results.append(advice)
    print(f"Product {i+1}: {advice}")

prs=Presentation()
# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Tea Management System Analysis"
slide.placeholders[1].text = "Decision Tree Classifier Results"

# Dataset Slide
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Dataset Overview"
rows, cols = df.shape
table = slide.shapes.add_table(rows+1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(2.5)).table

# Add headers
for j, col_name in enumerate(df.columns):
    table.cell(0, j).text = col_name

# Add rows
for i in range(rows):
    for j, col_name in enumerate(df.columns):
        table.cell(i+1, j).text = str(df.iloc[i, j])

# Decision Tree Slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Decision Tree Visualization"
slide.shapes.add_picture(tree_path, Inches(1), Inches(1.5), width=Inches(7))

# Feature Importance Slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Feature Importance"
slide.shapes.add_picture(feature_importances_path, Inches(1), Inches(1.5), width=Inches(7))

# Product Type vs Quantity Sold Slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Product Type vs Quantity Sold"
slide.shapes.add_picture(path, Inches(1), Inches(1.5), width=Inches(7))

# Save PPT
prs.save("Tea_Management_AI.pptx")
print("Presentation created successfully!")



























































